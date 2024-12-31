import sys
import docx
import PyPDF2
from pptx import Presentation
from os import listdir
from os.path import isfile, join, isdir
from langchain_text_splitters import TokenTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams
from langchain_qdrant import QdrantVectorStore

def lista_arquivos(dir):

    arquivos_list = []
    for f in listdir(dir):
        if isfile(join(dir, f)):
            arquivos_list.append(join(dir, f))
        elif isdir(join(dir, f)):
            arquivos_list += lista_arquivos(join(dir, f))
    
    return arquivos_list

def carrega_texto_word(arquivoname):

    doc = docx.Document(arquivoname)
    fullText = [para.text for para in doc.paragraphs]

    return '\n'.join(fullText)

def carrega_texto_pptx(arquivoname):

    prs = Presentation(arquivoname)
    fullText = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                fullText.append(shape.text)
                
    return '\n'.join(fullText)

def main_indexing(mypath):

    model_name = 'sentence-transformers/msmarco-bert-base-dot-v5'
    model_kwargs = {'device':'cpu'}
    encode_kwargs = {'normalize_embeddings': True}

    hf = HuggingFaceEmbeddings(model_name = model_name, model_kwargs= model_kwargs, encode_kwargs= encode_kwargs)

    client = QdrantClient('http://localhost:6333')

    collection_name = 'VectorDB'

    if client.collection_exists(collection_name):
        client.delete_collection(collection_name)

    client.create_collection(collection_name= collection_name, vectors_config= VectorParams(size= 768, distance= Distance.COSINE))
    qdrant = QdrantVectorStore(client= client, collection_name= collection_name, embedding= hf)
    print('\nIndexando os documentos...\n')

    arquivos = lista_arquivos(mypath)

    for arquivo in arquivos:

        try:
            arquivo_content = ''

            if arquivo.endswith('.pdf'):
                print(f'Indexando: {arquivo}')
                reader = PyPDF2.PdfReader(arquivo)
                for page in reader.pages:
                    arquivo_content += ' ' + page.extract_text()

            elif arquivo.endswith('.txt'):
                print(f'Indexando: {arquivo}')
                with open(arquivo, 'r') as f:
                    arquivo_content = f.read()

            elif arquivo.endswith('.docx'):
                print(f'Indexando: {arquivo}')
                arquivo_content = carrega_texto_word(arquivo)

            elif arquivo.endswith('.pptx'):
                print(f'Indexando: {arquivo}')
                arquivo_content = carrega_texto_pptx(arquivo)

            else:
                continue

            text_spliter = TokenTextSplitter(chunk_size = 500, chunk_overlap = 50)
            textos = text_spliter.split_text(arquivo_content)
            metadata = [{'path': arquivo} for _ in textos]
            qdrant.add_texts(textos, metadatas= metadata)

        except Exception as e:
            print(f'O processo falhou para o aarquivo {arquivo}: {e}')
    
    print('\nIndexação Concluída\n')

if __name__ == '__main__':
    arguments = sys.argv

    if len(arguments) > 1:
        main_indexing(arguments[1])
    else:
        print('Você precisa fornecer um caminho para a pasta com documentos para indexar.')